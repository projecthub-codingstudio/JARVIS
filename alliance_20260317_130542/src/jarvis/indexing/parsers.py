"""Document parsers — extract text content from various file formats.

Parser Tiers per PHASE1_ARCHITECTURE_CORE_DESIGN.md Section 12:
  Tier 1 (committed): Markdown, plain text, source code, PDF, DOCX, XLSX, PPTX
  Tier 2 (conditional): HWPX (formatting-loss caveats accepted)

Parser routing per Implementation Spec Task 1.1.

Text file auto-detection: files with unregistered extensions are probed
for text content and indexed as plain text if readable.
"""

from __future__ import annotations

import hashlib
import logging
import re
from pathlib import Path

from jarvis.contracts import AccessStatus, DocumentRecord, IndexingStatus

logger = logging.getLogger(__name__)

# --- Extension → type mapping ---
# Known extensions are mapped to a parser type.
# Unknown extensions go through text auto-detection (see is_indexable).

_EXTENSION_TYPE_MAP: dict[str, str] = {
    # Text / Markdown
    ".md": "markdown",
    ".markdown": "markdown",
    ".txt": "text",
    ".rst": "text",
    ".log": "text",
    ".csv": "text",
    ".tsv": "text",
    # Config / Data
    ".cfg": "text",
    ".toml": "text",
    ".ini": "text",
    ".env": "text",
    ".properties": "text",
    ".conf": "text",
    ".xml": "text",
    ".html": "text",
    ".htm": "text",
    ".css": "text",
    ".scss": "text",
    ".less": "text",
    ".svg": "text",
    # Code: Python
    ".py": "python",
    ".pyi": "python",
    # Code: JavaScript / TypeScript
    ".ts": "typescript",
    ".tsx": "typescript",
    ".js": "javascript",
    ".jsx": "javascript",
    # Code: JVM
    ".java": "code",
    ".kt": "code",
    ".kts": "code",
    ".scala": "code",
    ".groovy": "code",
    ".gradle": "code",
    # Code: C / C++ / Objective-C
    ".c": "code",
    ".h": "code",
    ".cpp": "code",
    ".hpp": "code",
    ".cc": "code",
    ".cxx": "code",
    ".m": "code",
    ".mm": "code",
    # Code: .NET
    ".cs": "code",
    ".fs": "code",
    ".vb": "code",
    # Code: Systems / Modern
    ".go": "code",
    ".rs": "code",
    ".swift": "code",
    ".zig": "code",
    # Code: Scripting
    ".rb": "code",
    ".php": "code",
    ".pl": "code",
    ".pm": "code",
    ".lua": "code",
    ".r": "code",
    ".jl": "code",
    ".ex": "code",
    ".exs": "code",
    ".erl": "code",
    ".clj": "code",
    ".cljs": "code",
    # Code: Shell
    ".sh": "code",
    ".bash": "code",
    ".zsh": "code",
    ".fish": "code",
    ".ps1": "code",
    ".bat": "code",
    ".cmd": "code",
    # Code: Config as code
    ".tf": "code",
    ".hcl": "code",
    ".proto": "code",
    ".graphql": "code",
    ".gql": "code",
    # Data / Markup
    ".yaml": "yaml",
    ".yml": "yaml",
    ".json": "json",
    ".json5": "json",
    ".jsonl": "json",
    # Code: SQL
    ".sql": "sql",
    ".ddl": "sql",
    # Tier 1: binary document formats
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    # Tier 2: Korean document formats
    ".hwpx": "hwpx",
    # Tier 3: legacy binary HWP (experimental)
    ".hwp": "hwp",
}

# Types that require specialized parsers (not plain read_text)
_BINARY_TYPES: frozenset[str] = frozenset({"pdf", "docx", "pptx", "xlsx", "hwpx", "hwp"})

# Extensions that are always binary — never attempt text detection
_BINARY_EXTENSIONS: frozenset[str] = frozenset({
    # Images
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".ico", ".webp", ".tiff", ".tif", ".svg",
    # Audio / Video
    ".mp3", ".mp4", ".wav", ".flac", ".ogg", ".avi", ".mov", ".mkv", ".webm",
    # Archives
    ".zip", ".tar", ".gz", ".bz2", ".xz", ".7z", ".rar",
    # Executables / Libraries
    ".exe", ".dll", ".so", ".dylib", ".o", ".a", ".pyc", ".pyo", ".class", ".wasm",
    # Fonts
    ".ttf", ".otf", ".woff", ".woff2",
    # Other binary
    ".db", ".sqlite", ".sqlite3", ".bin", ".dat", ".parquet", ".avro",
    ".DS_Store",
})

# Max file size for text auto-detection (skip very large unknown files)
_MAX_TEXT_PROBE_BYTES = 1_048_576  # 1MB

# Encoding probe order for text detection and reading.
# Covers: macOS/Linux UTF-8, Windows 10/11 Korean (CP949/EUC-KR),
# BOM-identified UTF-16 LE/BE, and legacy Latin-1.
#
# Note: BOM-less UTF-16 is intentionally excluded from the fallback chain.
# Python's utf-16 decoder can falsely "succeed" on normal UTF-8 Korean text
# and poison the persisted search index with mojibake.
_TEXT_ENCODINGS = ("utf-8", "cp949", "euc-kr", "latin-1")

# BOM signatures for UTF encoding detection
_BOM_MAP: list[tuple[bytes, str]] = [
    (b"\xff\xfe", "utf-16-le"),   # Windows "유니코드" (Notepad default)
    (b"\xfe\xff", "utf-16-be"),
    (b"\xef\xbb\xbf", "utf-8-sig"),  # Windows "UTF-8" (with BOM)
]


def _detect_encoding(sample: bytes) -> str | None:
    """Detect text encoding from a byte sample.

    Checks BOM first, then tries common encodings.
    Returns the encoding name or None if all fail.
    """
    # 1. BOM detection (Windows text editors)
    for bom, enc in _BOM_MAP:
        if sample.startswith(bom):
            return enc

    # 2. Try encodings in order
    for enc in _TEXT_ENCODINGS:
        try:
            sample.decode(enc)
            return enc
        except UnicodeDecodeError as exc:
            # Sample-based detection can cut a multibyte character at the end.
            # If the only failure is an incomplete trailing sequence, accept
            # the encoding and let the caller decode the full file.
            if exc.end == len(sample) and exc.start >= max(0, len(sample) - 4):
                try:
                    sample[:exc.start].decode(enc)
                    return enc
                except (UnicodeDecodeError, UnicodeError):
                    pass
        except UnicodeError:
            continue

    return None


def is_text_file(path: Path) -> bool:
    """Probe whether a file is text-readable.

    Handles Windows text files (CP949, EUC-KR, UTF-16 LE with BOM)
    and macOS/Linux UTF-8. Korean encoding is detected without data loss.
    """
    if path.suffix.lower() in _BINARY_EXTENSIONS:
        return False
    try:
        size = path.stat().st_size
    except OSError:
        return False
    if size == 0 or size > _MAX_TEXT_PROBE_BYTES:
        return False
    try:
        sample = path.read_bytes()[:8192]

        # UTF-16 files contain null bytes normally — check BOM first
        for bom, _enc in _BOM_MAP:
            if sample.startswith(bom):
                return True

        # For non-BOM files, null bytes indicate binary
        if b"\x00" in sample:
            return False

        return _detect_encoding(sample) is not None
    except OSError:
        return False


def is_indexable(path: Path) -> bool:
    """Check if a file can be indexed — registered extension OR text auto-detect."""
    suffix = path.suffix.lower()
    if suffix in _EXTENSION_TYPE_MAP:
        return True
    if suffix in _BINARY_EXTENSIONS:
        return False
    if path.name.startswith(".") or path.name.startswith("~$"):
        return False
    return is_text_file(path)


# --- Individual format parsers ---


def _parse_pdf(path: Path) -> str:
    """Extract text from PDF using PyMuPDF.

    Limits to first 500KB of text to avoid excessive chunking
    on very large PDFs (e.g. 16MB C# reference = 2.8M chars).
    Returns empty string if pymupdf is not installed (graceful degradation).
    """
    try:
        import pymupdf
    except ImportError:
        logger.warning("pymupdf not installed — skipping PDF: %s", path.name)
        return ""

    _MAX_TEXT_CHARS = 500_000  # ~250K tokens, more than enough for RAG

    pages: list[str] = []
    total_chars = 0
    with pymupdf.open(str(path)) as doc:
        for page in doc:
            text = page.get_text()
            if text.strip():
                pages.append(text)
                total_chars += len(text)
                if total_chars > _MAX_TEXT_CHARS:
                    logger.info("PDF truncated at %d chars: %s", total_chars, path.name)
                    break
    return "\n\n".join(pages)


_MIN_BLOCK_CHARS = 200  # Minimum characters per merged text block
_MAX_PDF_TEXT_CHARS = 500_000  # ~250K tokens


def _parse_pdf_structured(path: Path) -> list:
    """Extract structured elements from PDF using PyMuPDF block-level data.

    Instead of page.get_text() which flattens everything into raw text,
    uses page.get_text("blocks") to get individual text blocks with
    positions. Merges small blocks on the same page to avoid the
    micro-chunk problem (e.g. 씨샵.pdf: 19,543 chunks avg 85 chars).

    Also attempts table extraction via page.find_tables() when available.
    """
    from jarvis.contracts import DocumentElement

    try:
        import pymupdf
    except ImportError:
        logger.warning("pymupdf not installed — skipping PDF: %s", path.name)
        return []

    elements: list[DocumentElement] = []
    total_chars = 0

    with pymupdf.open(str(path)) as doc:
        for page_num, page in enumerate(doc):
            if total_chars > _MAX_PDF_TEXT_CHARS:
                logger.info(
                    "PDF truncated at page %d (%d chars): %s",
                    page_num, total_chars, path.name,
                )
                break

            # --- Table extraction (PyMuPDF ≥ 1.23) ---
            table_rects: list = []
            try:
                tables = page.find_tables()
                for table in tables:
                    table_rects.append(table.bbox)
                    extracted = table.extract()
                    if not extracted or len(extracted) < 2:
                        continue
                    headers = tuple(str(c) if c else "" for c in extracted[0])
                    rows = tuple(
                        tuple(str(c) if c else "" for c in row)
                        for row in extracted[1:]
                        if any(str(c).strip() for c in row if c)
                    )
                    if headers or rows:
                        text = f"[Page {page_num + 1}] {' | '.join(headers)}" if headers else f"[Page {page_num + 1}]"
                        elements.append(DocumentElement(
                            element_type="table",
                            text=text,
                            metadata={
                                "headers": headers,
                                "rows": rows,
                                "sheet_name": f"Page {page_num + 1}",
                            },
                        ))
            except Exception:
                # find_tables() not available or failed — proceed with text only
                pass

            # --- Text block extraction ---
            blocks = page.get_text("blocks")  # (x0, y0, x1, y1, text, block_no, block_type)
            text_blocks: list[str] = []

            for block in blocks:
                # block_type 0 = text, 1 = image
                if block[6] != 0:
                    continue

                block_text = block[4].strip()
                if not block_text:
                    continue

                # Skip blocks that fall inside a detected table region
                bx0, by0, bx1, by1 = block[0], block[1], block[2], block[3]
                in_table = False
                for trect in table_rects:
                    tx0, ty0, tx1, ty1 = trect
                    if bx0 >= tx0 - 2 and by0 >= ty0 - 2 and bx1 <= tx1 + 2 and by1 <= ty1 + 2:
                        in_table = True
                        break
                if in_table:
                    continue

                text_blocks.append(block_text)

            # --- Merge small adjacent blocks ---
            merged = _merge_small_blocks(text_blocks, min_chars=_MIN_BLOCK_CHARS)

            for block_text in merged:
                total_chars += len(block_text)
                elements.append(DocumentElement(
                    element_type="text",
                    text=block_text,
                    metadata={"page": page_num + 1},
                ))

                if total_chars > _MAX_PDF_TEXT_CHARS:
                    break

    return elements


def _merge_small_blocks(blocks: list[str], *, min_chars: int) -> list[str]:
    """Merge consecutive small text blocks until each meets min_chars.

    Blocks that are already large enough are emitted as-is.
    The last accumulated group is always emitted regardless of size.
    """
    if not blocks:
        return []

    merged: list[str] = []
    buf: list[str] = []
    buf_len = 0

    for block in blocks:
        buf.append(block)
        buf_len += len(block)

        if buf_len >= min_chars:
            merged.append("\n\n".join(buf))
            buf = []
            buf_len = 0

    # Emit remaining buffer
    if buf:
        # Try to attach to previous merged block if it exists and result is reasonable
        if merged and buf_len < min_chars // 2:
            merged[-1] = merged[-1] + "\n\n" + "\n\n".join(buf)
        else:
            merged.append("\n\n".join(buf))

    return merged


def _parse_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx.

    Returns empty string if python-docx is not installed (graceful degradation).
    """
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed — skipping DOCX: %s", path.name)
        return ""

    doc = Document(str(path))
    parts: list[str] = []

    # Paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n\n".join(parts)


def _parse_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx.

    Extracts slide titles, body text, table cells, and notes.
    Each slide is separated by a slide header for chunker context.
    Returns empty string if python-pptx is not installed (graceful degradation).
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed — skipping PPTX: %s", path.name)
        return ""

    prs = Presentation(str(path))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_parts.append(" | ".join(cells))

        # Slide notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[Notes] {notes}")

        if slide_parts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))

    return "\n\n".join(parts)


def _parse_xlsx(path: Path) -> str:
    """Extract text from XLSX using openpyxl.

    Each row becomes a separate paragraph (\\n\\n separated)
    so the Chunker can split at row boundaries.
    Header row is prepended to each data row for context.
    Returns empty string if openpyxl is not installed (graceful degradation).
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed — skipping XLSX: %s", path.name)
        return ""

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        header: str = ""
        row_idx = 0

        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if not cells:
                continue

            row_text = " | ".join(cells)

            if row_idx == 0:
                # First row is header
                header = row_text
                parts.append(f"[{sheet_name}] {header}")
            else:
                # Data rows: include header mapping for context
                parts.append(f"[{sheet_name}] {row_text}")

            row_idx += 1

    wb.close()
    return "\n\n".join(parts)


def _parse_hwpx(path: Path) -> str:
    """Extract text from HWPX (Korean word processor, XML-based ZIP).

    HWPX is a ZIP archive containing OWPML XML (KS X 6101).
    Formatting loss is accepted per spec.
    """
    try:
        from hwpx import HWPXFile

        hwpx_file = HWPXFile(str(path))
        return hwpx_file.get_text()
    except Exception:
        # Fallback: manual ZIP+XML extraction
        return _parse_hwpx_fallback(path)


def _parse_hwpx_structured(path: Path) -> list:
    """Extract structured elements (text + tables) from HWPX.

    HWPX is a ZIP archive with OWPML XML (KS X 6101).
    Extracts text paragraphs and tables as separate DocumentElements.
    Uses hwpx library for text, and direct XML parsing for tables.
    """
    from jarvis.contracts import DocumentElement

    elements: list[DocumentElement] = []

    # 1. Extract text via hwpx library (best quality)
    try:
        from hwpx import HwpxPackage, TextExtractor

        pkg = HwpxPackage.open(str(path))
        extractor = TextExtractor(pkg)
        text = extractor.extract_text(paragraph_separator="\n\n", skip_empty=True)
        extractor.close()
        pkg.close()
        if text.strip():
            elements.append(DocumentElement(element_type="text", text=text.strip()))
    except Exception:
        # Fallback to manual text extraction
        text = _parse_hwpx_fallback(path)
        if text.strip():
            elements.append(DocumentElement(element_type="text", text=text.strip()))

    # 2. Extract tables from XML directly
    tables = _extract_hwpx_tables(path)
    elements.extend(tables)

    return elements


def _extract_hwpx_tables(path: Path) -> list:
    """Extract tables from HWPX ZIP by parsing XML for <hp:tbl> elements."""
    import zipfile
    from xml.etree import ElementTree
    from jarvis.contracts import DocumentElement

    _HP_NS = "http://www.hancom.co.kr/hwpml/2011/paragraph"
    tables: list[DocumentElement] = []
    table_idx = 0

    def _local_name(elem: ElementTree.Element) -> str:
        return elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

    def _normalized_text(elem: ElementTree.Element) -> str:
        return " ".join("".join(elem.itertext()).split())

    def _caption_from_neighbors(
        tbl_elem: ElementTree.Element,
        parent_map: dict[ElementTree.Element, ElementTree.Element],
    ) -> str:
        for elem in tbl_elem.iter():
            if "caption" in _local_name(elem).lower():
                caption = _normalized_text(elem)
                if caption:
                    return caption

        parent = parent_map.get(tbl_elem)
        if parent is None:
            return ""

        siblings = list(parent)
        try:
            index = siblings.index(tbl_elem)
        except ValueError:
            return ""

        nearby = siblings[max(0, index - 2): min(len(siblings), index + 2)]
        for sibling in nearby:
            if sibling is tbl_elem:
                continue
            sibling_name = _local_name(sibling).lower()
            sibling_text = _normalized_text(sibling)
            if not sibling_text:
                continue
            if "caption" in sibling_name:
                return sibling_text
            if sibling_name in {"p", "paragraph"} and re.match(r"^(표|table)\s*\d+", sibling_text, re.IGNORECASE):
                return sibling_text

        return ""

    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            for name in sorted(zf.namelist()):
                if not (name.startswith("Contents/") and name.endswith(".xml")):
                    continue
                try:
                    xml_bytes = zf.read(name)
                    root = ElementTree.fromstring(xml_bytes)
                    parent_map = {child: parent for parent in root.iter() for child in parent}
                except (ElementTree.ParseError, KeyError):
                    continue

                # Find all table elements (hp:tbl)
                for tbl_elem in root.iter(f"{{{_HP_NS}}}tbl"):
                    rows_data: list[list[str]] = []
                    for tr_elem in tbl_elem.iter(f"{{{_HP_NS}}}tr"):
                        cells: list[str] = []
                        for tc_elem in tr_elem.iter(f"{{{_HP_NS}}}tc"):
                            # Collect all text within this cell
                            cell_texts = []
                            for t_elem in tc_elem.iter(f"{{{_HP_NS}}}t"):
                                if t_elem.text:
                                    cell_texts.append(t_elem.text.strip())
                            cells.append(" ".join(cell_texts))
                        if cells:
                            rows_data.append(cells)

                    if len(rows_data) < 2:
                        continue

                    table_idx += 1
                    headers = tuple(rows_data[0])
                    rows = tuple(tuple(r) for r in rows_data[1:] if any(r))
                    caption = _caption_from_neighbors(tbl_elem, parent_map)
                    label = caption if caption else f"Table {table_idx}"
                    text_repr = f"[{label}] {' | '.join(headers)}" if headers else f"[{label}]"
                    tables.append(DocumentElement(
                        element_type="table",
                        text=text_repr,
                        metadata={
                            "headers": headers,
                            "rows": rows,
                            "sheet_name": label,
                        },
                    ))
    except Exception as exc:
        logger.warning("HWPX table extraction failed: %s", exc)

    return tables


def _parse_hwpx_fallback(path: Path) -> str:
    """Fallback HWPX parser using direct ZIP+XML extraction."""
    import zipfile
    from xml.etree import ElementTree

    parts: list[str] = []

    with zipfile.ZipFile(str(path), "r") as zf:
        for name in zf.namelist():
            if name.startswith("Contents/") and name.endswith(".xml"):
                try:
                    xml_bytes = zf.read(name)
                    root = ElementTree.fromstring(xml_bytes)
                    # Extract all text content from XML
                    for elem in root.iter():
                        if elem.text and elem.text.strip():
                            parts.append(elem.text.strip())
                        if elem.tail and elem.tail.strip():
                            parts.append(elem.tail.strip())
                except ElementTree.ParseError:
                    continue

    return "\n".join(parts)


def _parse_sql(path: Path) -> str:
    """Extract structured text from SQL files with encoding detection.

    Handles legacy Korean SQL files (EUC-KR, CP949) commonly exported from
    MS SQL Server Management Studio and Korean database tools.

    Produces a structured output optimized for RAG retrieval:
      1. Summary header with table/view names
      2. Column definitions table (name | type | nullable | description)
      3. Consolidated SQL body (GO statements collapsed)
    """
    import re as _re

    raw = path.read_bytes()

    # Encoding detection via shared BOM-aware detector
    detected = _detect_encoding(raw[:8192])
    if detected:
        try:
            text = raw.decode(detected)
        except (UnicodeDecodeError, UnicodeError):
            text = ""
    else:
        text = ""

    if not text:
        for encoding in _TEXT_ENCODINGS:
            try:
                text = raw.decode(encoding)
                break
            except (UnicodeDecodeError, UnicodeError):
                continue
    if not text:
        text = raw.decode("utf-8", errors="replace")

    # --- Extract structured metadata ---

    # Table names
    tables = _re.findall(
        r"CREATE\s+TABLE\s+(?:IF\s+NOT\s+EXISTS\s+)?(?:\[?(\w+)\]?\.)?\[?(\w+)\]?",
        text, _re.IGNORECASE,
    )
    # View names
    views = _re.findall(
        r"CREATE\s+(?:OR\s+REPLACE\s+)?VIEW\s+(?:\[?(\w+)\]?\.)?\[?(\w+)\]?",
        text, _re.IGNORECASE,
    )
    # Index names
    indexes = _re.findall(
        r"CREATE\s+(?:UNIQUE\s+)?(?:CLUSTERED\s+)?(?:NONCLUSTERED\s+)?INDEX\s+\[?(\w+)\]?",
        text, _re.IGNORECASE,
    )

    # Extended property descriptions: column_name → description
    # Each EXEC sp_addextendedproperty is a single line/block
    ext_prop_map: dict[str, str] = {}
    table_descs: dict[str, str] = {}
    for m in _re.finditer(
        r"sp_addextendedproperty\s+.*?@value\s*=\s*N'([^']+)'.*?@level1name\s*=\s*N'([^']+)'(.*?)(?:GO|$)",
        text, _re.IGNORECASE | _re.DOTALL,
    ):
        desc_value = m.group(1)
        table_name = m.group(2)
        remainder = m.group(3)
        col_match = _re.search(r"@level2name\s*=\s*N'([^']+)'", remainder, _re.IGNORECASE)
        if col_match:
            ext_prop_map[col_match.group(1)] = desc_value
        else:
            table_descs[table_name] = desc_value

    # --- Extract column definitions from CREATE TABLE ---
    column_sections: list[str] = []
    for ct_match in _re.finditer(
        r"CREATE\s+TABLE\s+(?:\[?\w+\]?\.)?\[?(\w+)\]?\s*\(",
        text, _re.IGNORECASE,
    ):
        table_name = ct_match.group(1)
        # Find matching closing paren using bracket counting
        start = ct_match.end() - 1  # position of opening '('
        depth = 0
        body_end = start
        for i in range(start, len(text)):
            if text[i] == "(":
                depth += 1
            elif text[i] == ")":
                depth -= 1
                if depth == 0:
                    body_end = i
                    break
        body = text[start + 1:body_end]

        # Parse column lines (before CONSTRAINT keyword)
        col_body = _re.split(r"\bCONSTRAINT\b", body, maxsplit=1, flags=_re.IGNORECASE)[0]
        columns: list[str] = []
        for col_match in _re.finditer(
            r"\[(\w+)\]\s+\[?(\w+)\]?(?:\(([^)]+)\))?\s*(NOT\s+NULL|NULL)?",
            col_body,
        ):
            col_name = col_match.group(1)
            col_type = col_match.group(2)
            col_size = f"({col_match.group(3)})" if col_match.group(3) else ""
            nullable = col_match.group(4) or ""
            desc = ext_prop_map.get(col_name, "")
            columns.append(
                f"  {col_name} | {col_type}{col_size} | {nullable} | {desc}"
            )

        if columns:
            table_label = table_name
            if table_name in table_descs:
                table_label = f"{table_name} ({table_descs[table_name]})"
            section = f"[Table: {table_label}]\n"
            section += "  column | type | nullable | description\n"
            section += "  " + "-" * 60 + "\n"
            section += "\n".join(columns)
            column_sections.append(section)

    # --- Build output ---
    output_parts: list[str] = []

    # 1. Summary header
    header_parts: list[str] = []
    if tables:
        names = [t[1] if t[1] else t[0] for t in tables]
        header_parts.append(f"[SQL Tables: {', '.join(names)}]")
    if views:
        names = [v[1] if v[1] else v[0] for v in views]
        header_parts.append(f"[SQL Views: {', '.join(names)}]")
    if indexes:
        header_parts.append(f"[SQL Indexes: {', '.join(indexes)}]")
    if header_parts:
        output_parts.append("\n".join(header_parts))

    if table_descs or ext_prop_map:
        description_parts: list[str] = []
        for table_name, desc in sorted(table_descs.items()):
            description_parts.append(f"{table_name}: {desc}")
        for column_name, desc in sorted(ext_prop_map.items()):
            description_parts.append(f"{column_name}: {desc}")
        output_parts.append(f"[Descriptions: {'; '.join(description_parts)}]")

    # 2. Column definitions (structured, easy for LLM to parse)
    if column_sections:
        output_parts.extend(column_sections)

    # 3. Consolidated SQL body (collapse GO noise and deduplicate metadata)
    clean_sql = _re.sub(r"\r\n", "\n", text)
    clean_sql = _re.sub(r"\nGO\s*\n", "\n", clean_sql)
    # Remove sp_addextendedproperty blocks (already in structured header)
    clean_sql = _re.sub(
        r"EXEC\s+sys\.sp_addextendedproperty\b.*?(?=\nEXEC|\nSET|\nCREATE|\Z)",
        "", clean_sql, flags=_re.IGNORECASE | _re.DOTALL,
    )
    # Remove SET noise (ANSI_NULLS, QUOTED_IDENTIFIER, ANSI_PADDING)
    clean_sql = _re.sub(
        r"SET\s+(?:ANSI_NULLS|QUOTED_IDENTIFIER|ANSI_PADDING)\s+(?:ON|OFF)\s*\n",
        "", clean_sql, flags=_re.IGNORECASE,
    )
    clean_sql = _re.sub(r"\n{3,}", "\n\n", clean_sql)
    clean_sql = clean_sql.strip()
    if clean_sql:
        output_parts.append(clean_sql)

    return "\n\n".join(output_parts)


def _resolve_hwp5_binary(name: str) -> str | None:
    """Find a pyhwp binary (hwp5txt, hwp5proc) in PATH or venv."""
    import shutil
    import sys
    binary = shutil.which(name)
    if binary is None:
        venv_bin = Path(sys.executable).parent / name
        if venv_bin.exists():
            binary = str(venv_bin)
    return binary


def _parse_hwp(path: Path) -> str:
    """Extract text + tables from legacy binary HWP.

    Uses hwp5proc xml → lxml parsing to preserve table content.
    Falls back to hwp5txt (text-only, tables lost) if XML parsing fails.
    Returns empty string if pyhwp is not available (graceful degradation).
    """
    import subprocess

    # Prefer hwp5proc xml: preserves table content as structured XML
    hwp5proc_bin = _resolve_hwp5_binary("hwp5proc")
    if hwp5proc_bin is not None:
        text = _parse_hwp_xml(hwp5proc_bin, path)
        if text:
            return text

    # Fallback: hwp5txt (text only, tables become <표> placeholders)
    hwp5txt_bin = _resolve_hwp5_binary("hwp5txt")
    if hwp5txt_bin is None:
        logger.warning("pyhwp not installed — skipping HWP: %s", path.name)
        return ""

    result = subprocess.run(
        [hwp5txt_bin, str(path)],
        capture_output=True, text=True, timeout=120, check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(f"hwp5txt failed: {result.stderr[:200]}")

    import re
    text = result.stdout
    text = re.sub(r"<그림>", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _parse_hwp_xml(hwp5proc_bin: str, path: Path) -> str:
    """Parse HWP via hwp5proc xml → lxml, extracting text and table content.

    Tables are rendered as pipe-delimited rows (same format as DOCX/XLSX tables)
    so the existing chunking pipeline (TableChunkStrategy) can process them.
    """
    import subprocess

    result = subprocess.run(
        [hwp5proc_bin, "xml", str(path)],
        capture_output=True, timeout=120, check=False,
    )
    if result.returncode != 0 or not result.stdout:
        return ""

    try:
        from lxml import etree

        root = etree.fromstring(result.stdout)
        parts: list[str] = []

        # Collect all TableBody/TableCell elements to skip their Text children
        # (table content is extracted via _extract_hwp_table_body instead)
        table_elements: set[int] = set()
        for tbl in root.iter():
            tbl_tag = tbl.tag.split("}")[-1] if "}" in tbl.tag else tbl.tag
            if tbl_tag in ("TableBody", "TableCell", "TableRow", "TableControl"):
                table_elements.add(id(tbl))
                for descendant in tbl.iter():
                    table_elements.add(id(descendant))

        # Walk all elements, extracting text and table content
        for elem in root.iter():
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

            if tag == "Text" and elem.text and elem.text.strip():
                # Skip text inside tables (already extracted by TableBody handler)
                if id(elem) not in table_elements:
                    parts.append(elem.text.strip())

            elif tag == "TableBody":
                table_text = _extract_hwp_table_body(elem)
                if table_text:
                    parts.append(table_text)

        import re
        text = "\n\n".join(parts)
        text = re.sub(r"<그림>", "", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    except Exception as exc:
        logger.warning("HWP XML parsing failed: %s — falling back to hwp5txt", exc)
        return ""


def _parse_hwp_structured(path: Path) -> list:
    """Parse HWP into separate text + table DocumentElements via hwp5proc XML.

    Each table becomes a DocumentElement(element_type="table") with headers/rows.
    Text between tables becomes DocumentElement(element_type="text").
    This enables the ChunkRouter to apply TableChunkStrategy for tables
    and ParagraphChunkStrategy for text.
    """
    import subprocess

    hwp5proc_bin = _resolve_hwp5_binary("hwp5proc")
    if hwp5proc_bin is None:
        return []

    result = subprocess.run(
        [hwp5proc_bin, "xml", str(path)],
        capture_output=True, timeout=120, check=False,
    )
    if result.returncode != 0 or not result.stdout:
        return []

    return _parse_hwp_structured_xml_bytes(result.stdout, path_name=path.name)


def _parse_hwp_structured_xml_bytes(xml_bytes: bytes, *, path_name: str = "") -> list:
    """Parse HWP XML bytes into structured text/table elements.

    This helper exists so HWP parser behavior can be tested without requiring
    `hwp5proc` in the test environment.
    """
    if not xml_bytes:
        return []

    try:
        from lxml import etree
        from jarvis.contracts import DocumentElement

        root = etree.fromstring(xml_bytes)
        elements: list[DocumentElement] = []

        # Build set of elements inside tables (to skip their Text)
        table_elements: set[object] = set()
        for tbl in root.iter():
            tbl_tag = tbl.tag.split("}")[-1] if "}" in tbl.tag else tbl.tag
            if tbl_tag in ("TableBody", "TableCell", "TableRow", "TableControl", "TableCaption"):
                table_elements.add(tbl)
                for desc in tbl.iter():
                    table_elements.add(desc)

        # Collect text paragraphs and tables in document order
        text_buffer: list[str] = []
        table_idx = 0
        consumed_text_elements: set[object] = set()
        current_heading_path = ""
        pending_heading_path = ""
        pending_parent_heading = ""

        def _flush_text():
            nonlocal text_buffer, current_heading_path, pending_heading_path, pending_parent_heading
            if text_buffer:
                for paragraph_text in text_buffer:
                    normalized_paragraph = paragraph_text.strip()
                    if not normalized_paragraph:
                        continue
                    compact_paragraph = " ".join(normalized_paragraph.split()).strip(" -:")
                    if _looks_like_hwp_parent_heading(compact_paragraph):
                        pending_parent_heading = compact_paragraph
                        pending_heading_path = ""
                        continue
                    if pending_parent_heading and compact_paragraph in _HWP_SUBHEADINGS:
                        current_heading_path = f"{pending_parent_heading} > {compact_paragraph}"
                        pending_parent_heading = ""
                        pending_heading_path = ""
                        continue
                    built_elements = _build_hwp_text_elements(normalized_paragraph)
                    paragraph_emitted = False
                    for built in built_elements:
                        heading_path = str(built.metadata.get("heading_path", "")).strip()
                        if heading_path:
                            current_heading_path = heading_path
                            pending_heading_path = ""
                            pending_parent_heading = ""
                            elements.append(built)
                            paragraph_emitted = True
                            continue
                        if pending_heading_path:
                            current_heading_path = pending_heading_path
                            pending_heading_path = ""
                        elif current_heading_path and _looks_like_hwp_followup_heading(built.text):
                            pending_heading_path = _derive_hwp_followup_heading(built.text, current_heading_path)
                            pending_parent_heading = ""
                            paragraph_emitted = True
                            break

                        if current_heading_path:
                            elements.append(DocumentElement(
                                element_type=built.element_type,
                                text=built.text,
                                metadata={**built.metadata, "heading_path": current_heading_path},
                            ))
                        else:
                            elements.append(built)
                        paragraph_emitted = True
                    if not paragraph_emitted and pending_heading_path:
                        continue
                text_buffer = []

        def _get_table_caption(body_elem) -> str:
            """Walk up to parent TableControl and find sibling TableCaption."""
            parent = body_elem.getparent()
            if parent is None:
                return ""
            ptag = parent.tag.split("}")[-1] if "}" in parent.tag else parent.tag
            if ptag != "TableControl":
                return ""
            for sibling in parent:
                stag = sibling.tag.split("}")[-1] if "}" in sibling.tag else sibling.tag
                if stag == "TableCaption":
                    cap = etree.tostring(sibling, method="text", encoding="unicode").strip()
                    return " ".join(cap.split())
            return ""

        for elem in root.iter():
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

            if tag in ("P", "Paragraph") and elem not in table_elements:
                paragraph_text = etree.tostring(elem, method="text", encoding="unicode").strip()
                paragraph_text = " ".join(paragraph_text.split())
                if paragraph_text:
                    text_buffer.append(paragraph_text)
                    for desc in elem.iter():
                        consumed_text_elements.add(desc)

            elif tag == "Text" and elem.text and elem.text.strip():
                if elem not in table_elements:
                    if elem in consumed_text_elements:
                        continue
                    text_buffer.append(elem.text.strip())

            elif tag == "TableBody":
                _flush_text()

                rows_data = _extract_hwp_table_rows(elem)
                if rows_data and len(rows_data) >= 2:
                    table_idx += 1
                    headers = tuple(rows_data[0])
                    rows = tuple(tuple(r) for r in rows_data[1:] if any(r))
                    # Use document caption (e.g., "표 29 그러데이션 유형")
                    caption = _get_table_caption(elem)
                    label = caption if caption else f"HWP Table {table_idx}"
                    text_repr = f"[{label}] {' | '.join(headers)}" if headers else f"[{label}]"
                    elements.append(DocumentElement(
                        element_type="table",
                        text=text_repr,
                        metadata={
                            "headers": headers,
                            "rows": rows,
                            "sheet_name": label,
                        },
                    ))
                elif rows_data:
                    text_buffer.append(" | ".join(rows_data[0]))

        _flush_text()

        logger.info("HWP structured parse: %d elements (%d tables) from %s",
                     len(elements), table_idx, path_name)
        return elements

    except Exception as exc:
        logger.warning("HWP structured parse failed: %s", exc)
        return []


_HWP_SUBHEADINGS = ("기본 구조", "저장 구조", "공통 속성", "개체 속성", "텍스트 정보", "텍스트 속성")
_HWP_FOLLOWUP_HEADING_EXCLUDE_PREFIXES = ("표 ", "[표", "그림 ", "[그림")
_HWP_EXPLANATORY_RE = re.compile(r"(?:이다|있다|한다|된다|저장된다|설명|구조)")


def _build_hwp_text_elements(text: str) -> list:
    """Promote inline section titles in HWP spec text into heading-aware elements."""
    from jarvis.contracts import DocumentElement

    normalized = " ".join(text.split())
    head1 = ""
    head2 = ""
    body = normalized
    for subheading in _HWP_SUBHEADINGS:
        marker = f" {subheading} "
        if marker not in normalized:
            continue
        prefix, suffix = normalized.split(marker, 1)
        if not prefix.strip() or not suffix.strip():
            continue
        head1_candidate = prefix.strip().split(".")[-1].strip()
        if len(head1_candidate) < 4 or len(head1_candidate) > 40:
            continue
        if not any(keyword in head1_candidate for keyword in ("자료 구조", "파일 구조", "공통 속성", "개체")):
            continue
        head1 = head1_candidate
        head2 = subheading
        body = suffix.strip()
        break

    if not head1:
        return [DocumentElement(element_type="text", text=text)]

    heading_path = f"{head1} > {head2}" if head2 else head1
    return [
        DocumentElement(
            element_type="text",
            text=body,
            metadata={"heading_path": heading_path},
        )
    ]


def _looks_like_hwp_followup_heading(text: str) -> bool:
    normalized = " ".join(text.split()).strip(" -:")
    if len(normalized) < 4 or len(normalized) > 40:
        return False
    if any(normalized.startswith(prefix) for prefix in _HWP_FOLLOWUP_HEADING_EXCLUDE_PREFIXES):
        return False
    if _HWP_EXPLANATORY_RE.search(normalized):
        return False
    if any(ch in normalized for ch in ".!?=|"):
        return False
    return True


def _looks_like_hwp_parent_heading(text: str) -> bool:
    normalized = " ".join(text.split()).strip(" -:")
    if len(normalized) < 4 or len(normalized) > 40:
        return False
    if any(ch in normalized for ch in ".!?=|"):
        return False
    if any(token in normalized for token in ("이다", "있다", "한다", "된다", "설명")):
        return False
    return any(keyword in normalized for keyword in ("자료 구조", "파일 구조", "공통 속성", "개체 정보"))


def _derive_hwp_followup_heading(text: str, current_heading_path: str) -> str:
    normalized = " ".join(text.split()).strip(" -:")
    if not normalized or not current_heading_path:
        return ""
    if " > " in current_heading_path:
        parent = current_heading_path.rsplit(" > ", 1)[0]
    else:
        parent = current_heading_path
    return f"{parent} > {normalized}" if parent else normalized


def _extract_hwp_table_rows(table_body) -> list[list[str]]:
    """Extract rows from a TableBody element as list of cell lists."""
    from lxml import etree

    rows: list[list[str]] = []
    for row_elem in table_body:
        tag = row_elem.tag.split("}")[-1] if "}" in row_elem.tag else row_elem.tag
        if tag != "TableRow":
            continue
        cells: list[str] = []
        for cell_elem in row_elem:
            cell_tag = cell_elem.tag.split("}")[-1] if "}" in cell_elem.tag else cell_elem.tag
            if cell_tag != "TableCell":
                continue
            cell_text = etree.tostring(cell_elem, method="text", encoding="unicode").strip()
            cell_text = " ".join(cell_text.split())
            cells.append(cell_text)
        if any(cells):
            rows.append(cells)
    return rows


def _extract_hwp_table_body(table_body) -> str:
    """Extract rows from a TableBody element as pipe-delimited text."""
    from lxml import etree

    rows: list[str] = []
    for row_elem in table_body:
        tag = row_elem.tag.split("}")[-1] if "}" in row_elem.tag else row_elem.tag
        if tag != "TableRow":
            continue
        cells: list[str] = []
        for cell_elem in row_elem:
            cell_tag = cell_elem.tag.split("}")[-1] if "}" in cell_elem.tag else cell_elem.tag
            if cell_tag != "TableCell":
                continue
            cell_text = etree.tostring(cell_elem, method="text", encoding="unicode").strip()
            cell_text = " ".join(cell_text.split())  # Normalize whitespace
            cells.append(cell_text)
        if any(cells):
            rows.append(" | ".join(cells))

    return "\n".join(rows) if rows else ""


# --- Parser dispatch ---

_PARSER_DISPATCH: dict[str, object] = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "xlsx": _parse_xlsx,
    "hwpx": _parse_hwpx,
    "hwp": _parse_hwp,
    "sql": _parse_sql,
}

# Types that use specialized parsers (binary formats + encoding-sensitive text formats)
_DISPATCHED_TYPES: frozenset[str] = _BINARY_TYPES | frozenset({"sql"})


class DocumentParser:
    """Parses files into raw text for downstream chunking and indexing.

    Routes to specialized parsers for binary formats (PDF, DOCX, XLSX, PPTX, HWPX, HWP).
    Falls back to UTF-8 read_text for text-based formats.
    Supports text auto-detection for unregistered extensions.
    """

    def detect_type(self, path: Path) -> str:
        """Detect the document type from file extension.

        Returns the mapped type for known extensions, or "text" for
        unregistered extensions that pass text auto-detection.
        """
        suffix = Path(path).suffix.lower()
        return _EXTENSION_TYPE_MAP.get(suffix, "text")

    def parse(self, path: Path) -> str:
        """Extract text content from a file.

        Routes to specialized parsers for binary formats.
        Falls back to encoding-aware text read for all other formats.

        Raises:
            FileNotFoundError: If the file does not exist.
        """
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"File not found: {path}")

        doc_type = self.detect_type(p)

        if doc_type in _DISPATCHED_TYPES:
            parser_fn = _PARSER_DISPATCH[doc_type]
            try:
                return parser_fn(p)  # type: ignore[operator]
            except ImportError:
                # Dependency not installed — already logged inside parser
                return ""
            except Exception as e:
                logger.warning("Parser failed for %s (%s): %s", p, doc_type, e)
                raise

        return self._read_text_with_fallback(p)

    def _read_text_with_fallback(self, path: Path) -> str:
        """Read text with BOM-aware encoding detection.

        Handles Windows 10/11 Korean text files (CP949, EUC-KR),
        Windows Notepad UTF-16 LE/BE, and UTF-8 with BOM.
        """
        raw = path.read_bytes()

        # BOM detection first (most reliable for Windows files)
        detected = _detect_encoding(raw[:8192])
        if detected:
            try:
                return raw.decode(detected)
            except (UnicodeDecodeError, UnicodeError):
                pass

        # Full fallback chain
        for encoding in _TEXT_ENCODINGS:
            try:
                return raw.decode(encoding)
            except (UnicodeDecodeError, UnicodeError):
                continue
        return raw.decode("utf-8", errors="replace")

    def parse_structured(self, path: Path) -> "ParsedDocument":
        """Parse a file into a structured ParsedDocument with typed elements.

        Routes to type-specific element extraction. Falls back to
        wrapping plain text as a single TextElement.
        """
        from jarvis.contracts import DocumentElement, ParsedDocument

        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"File not found: {path}")

        doc_type = self.detect_type(p)
        suffix = p.suffix.lower()
        metadata = {"filename": p.name, "format": doc_type}

        # xlsx/csv → TableElement with headers and rows
        if doc_type == "xlsx":
            elements = self._parse_xlsx_structured(p)
            return ParsedDocument(elements=tuple(elements), metadata=metadata)

        # Code files → CodeElement
        code_types = {
            "python", "typescript", "javascript", "code",
            "c", "cpp", "go", "rust", "java", "swift", "kotlin",
        }
        if doc_type in code_types:
            text = self.parse(p)
            if text.strip():
                lang_map = {"python": "python", "typescript": "typescript", "javascript": "javascript",
                            "swift": "swift", "java": "java", "go": "go", "rust": "rust"}
                lang = lang_map.get(doc_type, doc_type)
                return ParsedDocument(
                    elements=(DocumentElement(element_type="code", text=text, metadata={"language": lang}),),
                    metadata=metadata,
                )
            return ParsedDocument(elements=(), metadata=metadata)

        # csv/tsv → TableElement
        if suffix in (".csv", ".tsv"):
            elements = self._parse_csv_structured(p, delimiter="\t" if suffix == ".tsv" else ",")
            return ParsedDocument(elements=tuple(elements), metadata=metadata)

        # docx → text + table elements
        if doc_type == "docx":
            elements = self._parse_docx_structured(p)
            return ParsedDocument(elements=tuple(elements), metadata=metadata)

        # HWP (binary) → text + table elements via hwp5proc XML
        if doc_type == "hwp":
            elements = _parse_hwp_structured(p)
            if elements:
                return ParsedDocument(elements=tuple(elements), metadata=metadata)
            text = self.parse(p)
            if text.strip():
                return ParsedDocument(
                    elements=(DocumentElement(element_type="text", text=text),),
                    metadata=metadata,
                )
            return ParsedDocument(elements=(), metadata=metadata)

        # HWPX → text + table elements
        if doc_type == "hwpx":
            elements = _parse_hwpx_structured(p)
            if elements:
                return ParsedDocument(elements=tuple(elements), metadata=metadata)
            text = self.parse(p)
            if text.strip():
                return ParsedDocument(
                    elements=(DocumentElement(element_type="text", text=text),),
                    metadata=metadata,
                )
            return ParsedDocument(elements=(), metadata=metadata)

        # PDF → structured block-level extraction
        if doc_type == "pdf":
            elements = _parse_pdf_structured(p)
            if elements:
                return ParsedDocument(elements=tuple(elements), metadata=metadata)
            # Fallback to plain text if structured extraction returned nothing
            text = self.parse(p)
            if text.strip():
                return ParsedDocument(
                    elements=(DocumentElement(element_type="text", text=text),),
                    metadata=metadata,
                )
            return ParsedDocument(elements=(), metadata=metadata)

        # All other formats → plain text wrapped as TextElement
        text = self.parse(p)
        if not text.strip():
            return ParsedDocument(elements=(), metadata=metadata)

        return ParsedDocument(
            elements=(DocumentElement(element_type="text", text=text),),
            metadata=metadata,
        )

    def _parse_xlsx_structured(self, path: Path) -> list:
        """Parse xlsx into TableElement per sheet."""
        from jarvis.contracts import DocumentElement
        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.warning("openpyxl not installed — skipping XLSX: %s", path.name)
            return []

        elements = []
        wb = load_workbook(str(path), read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            headers: tuple[str, ...] = ()
            rows: list[tuple[str, ...]] = []
            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                cells = tuple(str(c) if c is not None else "" for c in row)
                if not any(cells):
                    continue
                if row_idx == 0:
                    headers = cells
                else:
                    rows.append(cells)
            if headers or rows:
                text = f"[{sheet_name}] {' | '.join(headers)}" if headers else f"[{sheet_name}]"
                elements.append(DocumentElement(
                    element_type="table",
                    text=text,
                    metadata={
                        "headers": headers,
                        "rows": tuple(rows),
                        "sheet_name": sheet_name,
                    },
                ))
        wb.close()
        return elements

    def _parse_csv_structured(self, path: Path, *, delimiter: str = ",") -> list:
        """Parse CSV/TSV into TableElement."""
        from jarvis.contracts import DocumentElement
        text = self._read_text_with_fallback(path)
        if not text.strip():
            return []

        lines = [l for l in text.strip().split("\n") if l.strip()]
        if not lines:
            return []

        headers = tuple(cell.strip() for cell in lines[0].split(delimiter))
        rows = tuple(
            tuple(cell.strip() for cell in line.split(delimiter))
            for line in lines[1:]
            if line.strip()
        )
        return [DocumentElement(
            element_type="table",
            text=f"[{path.stem}] {' | '.join(headers)}",
            metadata={"headers": headers, "rows": rows, "sheet_name": path.stem},
        )]

    def _parse_docx_structured(self, path: Path) -> list:
        """Parse DOCX into text + table elements."""
        from jarvis.contracts import DocumentElement
        try:
            from docx import Document
        except ImportError:
            logger.warning("python-docx not installed — skipping DOCX: %s", path.name)
            return []

        elements = []
        doc = Document(str(path))

        # Paragraphs as text
        text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
        if text_parts:
            elements.append(DocumentElement(
                element_type="text",
                text="\n\n".join(text_parts),
            ))

        # Tables as TableElement
        for table_idx, table in enumerate(doc.tables):
            headers: tuple[str, ...] = ()
            rows: list[tuple[str, ...]] = []
            for row_idx, row in enumerate(table.rows):
                cells = tuple(cell.text.strip() for cell in row.cells)
                if not any(cells):
                    continue
                if row_idx == 0:
                    headers = cells
                else:
                    rows.append(cells)
            if headers or rows:
                elements.append(DocumentElement(
                    element_type="table",
                    text=f"[Table {table_idx + 1}] {' | '.join(headers)}" if headers else f"[Table {table_idx + 1}]",
                    metadata={"headers": headers, "rows": tuple(rows), "sheet_name": f"Table {table_idx + 1}"},
                ))
        return elements

    def create_record(self, path: Path) -> DocumentRecord:
        """Create a DocumentRecord for a file.

        Raises:
            FileNotFoundError: If the file does not exist.
        """
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"File not found: {path}")

        raw_bytes = p.read_bytes()
        content_hash = hashlib.sha256(raw_bytes).hexdigest()
        size_bytes = p.stat().st_size

        return DocumentRecord(
            path=str(p),
            content_hash=content_hash,
            size_bytes=size_bytes,
            indexing_status=IndexingStatus.PENDING,
            access_status=AccessStatus.ACCESSIBLE,
        )
